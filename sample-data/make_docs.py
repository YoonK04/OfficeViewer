import os
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPt
from pptx.dml.color import RGBColor as PColor

base = os.path.dirname(os.path.abspath(__file__))

# ── Word ─────────────────────────────
doc = Document()
doc.add_heading("프로젝트 진행 보고서", level=0)
doc.add_heading("1. 개요", level=1)
p = doc.add_paragraph("이 문서는 Office Viewer의 ")
r = p.add_run("워드 미리보기")
r.bold = True
r.font.color.rgb = RGBColor(0x0E, 0x63, 0x9C)
p.add_run(" 기능을 테스트하기 위한 샘플입니다.")
doc.add_heading("2. 항목", level=1)
for item in ["일정 점검", "예산 검토", "리스크 관리"]:
    doc.add_paragraph(item, style="List Bullet")
doc.add_heading("3. 표", level=1)
table = doc.add_table(rows=3, cols=3)
table.style = "Light Grid Accent 1"
hdr = table.rows[0].cells
for i, h in enumerate(["구분", "계획", "실적"]):
    hdr[i].text = h
data = [["1분기", "100", "98"], ["2분기", "120", "131"]]
for r_i, row in enumerate(data, start=1):
    for c_i, v in enumerate(row):
        table.rows[r_i].cells[c_i].text = v
doc.save(os.path.join(base, "보고서.docx"))

# ── PowerPoint ───────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Office Viewer 데모"
slide.placeholders[1].text = "통합 오피스 뷰어 — PPT 미리보기 테스트"

slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "주요 기능"
body = slide2.placeholders[1].text_frame
body.text = "파일 탐색기"
for t in ["탭으로 열기", "좌우 분할 보기", "엑셀 보기+편집"]:
    para = body.add_paragraph()
    para.text = t

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
box = slide3.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
tf = box.text_frame
run = tf.paragraphs[0].add_run()
run.text = "감사합니다"
run.font.size = PPt(54)
run.font.bold = True
run.font.color.rgb = PColor(0x0E, 0x63, 0x9C)

prs.save(os.path.join(base, "발표자료.pptx"))
print("docx/pptx created:", os.listdir(base))
